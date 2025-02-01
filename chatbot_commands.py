import re
import db_commands as db
from difflib import get_close_matches
from fuzzywuzzy import fuzz
import fitz
import docx
import pptx
from io import BytesIO
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np
import pandas as pd
from pathlib import  Path


def find_table(user_input):
    user_input = user_input.lower().strip().replace(" ", "_")

    tables = db.list_bd()

    if not tables:
        return None

    matches = get_close_matches(user_input, tables, n=1, cutoff=0.4)
    return matches[0] if matches else None


def select_table(user_input):
    clossest_table = find_table(user_input)

    if clossest_table:
        return clossest_table, f"Tabela '{clossest_table}' selecionada"
    else:
        return None, f"Erro: não encontrei nenhuma tabela parecida"


def get_relevant_files(user_input, user_question):
    try:
        table_name, msg = select_table(user_input)
        if not table_name:
            return None, msg

        files = db.select_bd(table_name)

        # **🚨 Verifica se `files` é None ou está vazio**
        if files is None or (isinstance(files, pd.DataFrame) and files.empty):
            return None, "Nenhum arquivo encontrado nessa matéria."

        # **✅ Converte para DataFrame se necessário**
        if isinstance(files, list):
            files = pd.DataFrame(files)

        # **🛠 Verificação de colunas**
        required_columns = {"name_archive", "description_archive", "category", }
        if not required_columns.issubset(files.columns):
            return None, "Erro: Estrutura da tabela está incorreta."

        # Supondo que 'files' seja o DataFrame
        allowed_extensions = (".pdf", ".txt", ".pptx", ".docx")

        filtered_files = files[files["category"] == "teoria"]
        filtered_files = filtered_files[filtered_files["name_archive"].str.endswith(tuple(allowed_extensions), na=False)]

        # **🎯 Ranqueando arquivos**
        files_score = []
        for _, file in filtered_files.iterrows():
            name_score = fuzz.partial_ratio(user_question.lower(), file["name_archive"].lower())
            desc_score = fuzz.partial_ratio(user_question.lower(), (file["description_archive"] or "").lower())
            total_score = (name_score + desc_score) / 2
            files_score.append((file, total_score))

        files_score.sort(key=lambda x: x[1], reverse=True)
        selected_files = [file[0] for file in files_score[:3]]

        return selected_files, "Arquivo selecionado com sucesso!!!"

    except Exception as e:
        return None, f"Erro ao buscar arquivos: {str(e)}"


def read_file(file_data, file_type):
    try:
        if file_type == "txt":
            return file_data.decode("utf-8", errors="ignore")

        elif file_type == "pdf":
            text = ""
            pdf_document = fitz.open(stream=BytesIO(file_data), filetype="pdf")
            for page in pdf_document:
                text += page.get_text("text") + "\n"
            return text.strip()

        elif file_type == "docx":
            text = ""
            doc = docx.Document(BytesIO(file_data))
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()

        elif file_type == "pptx":
            text = ""
            presentation = pptx.Presentation(BytesIO(file_data))
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()

        else:
            return "Erro: Formato de arquivo não suportado."

    except Exception as e:
        return f"Erro ao ler o arquivo: {str(e)}"


def relevant_snippets(user_question, files_text, num_snippets=3):

    def split_text(text):
        return text.split("\n\n")

    snippets = []
    for file_text in files_text:
        snippets.extend(split_text(file_text))

    if not snippets:
        return ["Nenhum trecho relevante encontrado"]

    vectorizer = TfidfVectorizer()
    tfidf_matrix = vectorizer.fit_transform([user_question] + snippets)

    similarities = (tfidf_matrix * tfidf_matrix.T).toarray()[0, 1:]

    '''similarities = np.array(similarities)
    if similarities.size >= num_snippets:
        best_indices = np.argsort(similarities)[-num_snippets:][::-1]
    else:
        best_indices = np.argsort(similarities)[::-1]

    best_snippets = [snippets[i] for i in best_indices[:num_snippets]]'''

    best_indices = np.argsort(similarities)[-num_snippets:][::-1]
    best_snippets = [snippets[i] for i in best_indices if len(snippets[i]) > 50]  # Filtra trechos muito curtos

    if not best_snippets:
        return ["Nenhum trecho relevante encontrado"]

    return best_snippets


# Simulando uma pergunta do usuário
user_input = "banco de dados"
user_question = "como instalar o sql server"

# Passo 1: Encontrar a tabela mais relevante
selected_files, msg = get_relevant_files(user_input, user_question)

if selected_files:
    print("Arquivos mais relevantes encontrados:")
    files_text = []

    # Passo 2: Ler os arquivos relevantes
    for file in selected_files:
        file_data = file["archive"]
        file_type = Path(file["name_archive"]).suffix.lstrip(".")
        file_text = read_file(file_data, file_type)
        files_text.append(file_text)

    # Passo 3: Encontrar os trechos mais relevantes
    snippets = relevant_snippets(user_question, files_text)

    print("\nTrechos mais relevantes:")
    for snippet in snippets:
        print(f"- {snippet}")

else:
    print(msg)





